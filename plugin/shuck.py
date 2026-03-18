#!/usr/bin/env python3
"""
shuck — Convert documents to clean Markdown for AI agents and LLMs.

Supports: .docx, .pdf, .xlsx, .pptx, .csv
"""

__version__ = "1.0.0"

import sys
import json
import argparse
from pathlib import Path

SUPPORTED_FORMATS = {
    ".docx": "Microsoft Word",
    ".pdf":  "PDF",
    ".xlsx": "Excel Spreadsheet",
    ".pptx": "PowerPoint Presentation",
    ".csv":  "CSV (Comma-Separated Values)",
}


# ── Word (.docx) conversion ─────────────────────────────────────────────


def convert_docx(filepath: Path) -> str:
    try:
        from docx import Document
        from docx.oxml.ns import qn
    except ImportError:
        print("Missing dependency: python-docx\n  pip install python-docx", file=sys.stderr)
        sys.exit(1)

    doc = Document(str(filepath))
    lines: list[str] = []

    for element in doc.element.body:
        tag = element.tag.split("}")[-1]

        if tag == "p":
            lines.append(_convert_paragraph(element, doc))
        elif tag == "tbl":
            lines.append(_convert_table(element, doc))

    return "\n".join(lines).strip() + "\n"


def _run_text(run) -> str:
    """Convert a single run to Markdown with bold/italic."""
    from docx.oxml.ns import qn

    text = run.text or ""
    if not text.strip():
        return text

    rpr = run.find(qn("w:rPr"))
    bold = False
    italic = False
    if rpr is not None:
        bold = rpr.find(qn("w:b")) is not None and rpr.find(qn("w:b")).get(qn("w:val"), "true") != "false"
        italic = rpr.find(qn("w:i")) is not None and rpr.find(qn("w:i")).get(qn("w:val"), "true") != "false"

    if bold and italic:
        return f"***{text}***"
    elif bold:
        return f"**{text}**"
    elif italic:
        return f"*{text}*"
    return text


def _paragraph_text(element) -> str:
    """Extract full text from a paragraph element with inline formatting."""
    from docx.oxml.ns import qn

    parts = []
    for run in element.findall(qn("w:r")):
        parts.append(_run_text(run))
    return "".join(parts)


def _get_style_name(element, doc) -> str:
    """Get the style name of a paragraph element."""
    from docx.oxml.ns import qn

    ppr = element.find(qn("w:pPr"))
    if ppr is not None:
        style_el = ppr.find(qn("w:pStyle"))
        if style_el is not None:
            return style_el.get(qn("w:val"), "")
    return ""


def _get_numpr(element):
    """Get numbering properties (list info) from a paragraph."""
    from docx.oxml.ns import qn

    ppr = element.find(qn("w:pPr"))
    if ppr is not None:
        return ppr.find(qn("w:numPr"))
    return None


def _convert_paragraph(element, doc) -> str:
    from docx.oxml.ns import qn

    style = _get_style_name(element, doc).lower()
    text = _paragraph_text(element)

    # Headings
    for level in range(1, 7):
        if style in (f"heading{level}", f"heading {level}"):
            return f"\n{'#' * level} {text}\n"

    # Title / Subtitle
    if style == "title":
        return f"\n# {text}\n"
    if style == "subtitle":
        return f"\n## {text}\n"

    # Lists
    numpr = _get_numpr(element)
    if numpr is not None or "list" in style:
        indent_level = 0
        if numpr is not None:
            ilvl = numpr.find(qn("w:ilvl"))
            if ilvl is not None:
                indent_level = int(ilvl.get(qn("w:val"), "0"))
        prefix = "  " * indent_level + "- "
        return prefix + text

    # Normal paragraph
    if not text.strip():
        return ""
    return text + "\n"


def _convert_table(element, doc) -> str:
    from docx.oxml.ns import qn

    rows = element.findall(qn("w:tr"))
    if not rows:
        return ""

    table_data: list[list[str]] = []
    for tr in rows:
        cells = []
        for tc in tr.findall(qn("w:tc")):
            cell_texts = []
            for p in tc.findall(qn("w:p")):
                cell_texts.append(_paragraph_text(p).strip())
            cells.append(" ".join(cell_texts).replace("|", "\\|"))
        table_data.append(cells)

    if not table_data:
        return ""

    # Normalize column count
    max_cols = max(len(row) for row in table_data)
    for row in table_data:
        while len(row) < max_cols:
            row.append("")

    md_lines = ["\n"]
    # Header row
    md_lines.append("| " + " | ".join(table_data[0]) + " |")
    md_lines.append("| " + " | ".join(["---"] * max_cols) + " |")
    # Data rows
    for row in table_data[1:]:
        md_lines.append("| " + " | ".join(row) + " |")
    md_lines.append("")

    return "\n".join(md_lines)


# ── PDF conversion ───────────────────────────────────────────────────────


def convert_pdf(filepath: Path) -> str:
    try:
        import pdfplumber
    except ImportError:
        print("Missing dependency: pdfplumber\n  pip install pdfplumber", file=sys.stderr)
        sys.exit(1)

    lines: list[str] = []
    with pdfplumber.open(str(filepath)) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text:
                lines.append(text)
            # Page separator
            if i < len(pdf.pages) - 1:
                lines.append("\n---\n")

    return "\n\n".join(lines).strip() + "\n"


# ── Excel (.xlsx) conversion ────────────────────────────────────────────


def convert_xlsx(filepath: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        print("Missing dependency: openpyxl\n  pip install openpyxl", file=sys.stderr)
        sys.exit(1)

    wb = load_workbook(str(filepath), read_only=True, data_only=True)
    sections: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))

        # Trim trailing empty rows
        while rows and all(c is None for c in rows[-1]):
            rows.pop()
        if not rows:
            sections.append(f"## Sheet: {sheet_name}\n\n*Empty sheet*\n")
            continue

        # Trim trailing empty columns
        max_col = 0
        for row in rows:
            for i in range(len(row) - 1, -1, -1):
                if row[i] is not None:
                    max_col = max(max_col, i + 1)
                    break
        if max_col == 0:
            sections.append(f"## Sheet: {sheet_name}\n\n*Empty sheet*\n")
            continue

        trimmed = [row[:max_col] for row in rows]

        # Build markdown table
        def cell_str(v):
            if v is None:
                return ""
            return str(v).replace("|", "\\|").replace("\n", " ")

        md_lines = [f"## Sheet: {sheet_name}\n"]
        header = trimmed[0]
        md_lines.append("| " + " | ".join(cell_str(c) for c in header) + " |")
        md_lines.append("| " + " | ".join(["---"] * max_col) + " |")
        for row in trimmed[1:]:
            md_lines.append("| " + " | ".join(cell_str(c) for c in row) + " |")
        md_lines.append("")
        sections.append("\n".join(md_lines))

    wb.close()
    return "\n".join(sections).strip() + "\n"


# ── PowerPoint (.pptx) conversion ───────────────────────────────────────


def convert_pptx(filepath: Path) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("Missing dependency: python-pptx\n  pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(str(filepath))
    sections: list[str] = []

    for idx, slide in enumerate(prs.slides, 1):
        slide_lines = [f"## Slide {idx}"]

        # Title
        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_lines.append(f"\n### {slide.shapes.title.text.strip()}\n")

        # Content shapes
        for shape in slide.shapes:
            # Skip the title shape we already handled
            if shape.shape_id == getattr(slide.shapes.title, "shape_id", None):
                continue

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_lines.append(text)

            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.replace("|", "\\|").replace("\n", " ") for cell in row.cells]
                    rows.append(cells)
                if rows:
                    max_cols = max(len(r) for r in rows)
                    for r in rows:
                        while len(r) < max_cols:
                            r.append("")
                    slide_lines.append("")
                    slide_lines.append("| " + " | ".join(rows[0]) + " |")
                    slide_lines.append("| " + " | ".join(["---"] * max_cols) + " |")
                    for r in rows[1:]:
                        slide_lines.append("| " + " | ".join(r) + " |")
                    slide_lines.append("")

        # Notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_lines.append(f"\n> **Notes:** {notes}\n")

        sections.append("\n".join(slide_lines))

    return "\n\n".join(sections).strip() + "\n"


# ── CSV conversion ──────────────────────────────────────────────────────


def convert_csv(filepath: Path) -> str:
    import csv

    text = filepath.read_text(encoding="utf-8", errors="replace")

    # Detect delimiter
    try:
        dialect = csv.Sniffer().sniff(text[:8192])
        delimiter = dialect.delimiter
    except csv.Error:
        delimiter = ","

    reader = csv.reader(text.splitlines(), delimiter=delimiter)
    rows = [row for row in reader if any(cell.strip() for cell in row)]

    if not rows:
        return "*Empty CSV file*\n"

    max_cols = max(len(r) for r in rows)
    for r in rows:
        while len(r) < max_cols:
            r.append("")

    def cell_str(v):
        return v.replace("|", "\\|").replace("\n", " ")

    md_lines = []
    md_lines.append("| " + " | ".join(cell_str(c) for c in rows[0]) + " |")
    md_lines.append("| " + " | ".join(["---"] * max_cols) + " |")
    for row in rows[1:]:
        md_lines.append("| " + " | ".join(cell_str(c) for c in row) + " |")

    return "\n".join(md_lines) + "\n"


# ── Meta info (agent probing) ───────────────────────────────────────────


def get_meta(filepath: Path) -> dict:
    ext = filepath.suffix.lower()
    meta = {
        "file": filepath.name,
        "format": ext.lstrip("."),
        "size_bytes": filepath.stat().st_size,
    }

    if ext == ".xlsx":
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(filepath), read_only=True, data_only=True)
            meta["sheets"] = wb.sheetnames
            wb.close()
        except Exception:
            meta["sheets"] = []

    elif ext == ".pdf":
        try:
            import pdfplumber
            with pdfplumber.open(str(filepath)) as pdf:
                meta["pages"] = len(pdf.pages)
        except Exception:
            meta["pages"] = None

    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(filepath))
            meta["slides"] = len(prs.slides)
        except Exception:
            meta["slides"] = None

    elif ext == ".csv":
        import csv
        try:
            text = filepath.read_text(encoding="utf-8", errors="replace")
            reader = csv.reader(text.splitlines())
            meta["rows"] = sum(1 for _ in reader)
        except Exception:
            meta["rows"] = None

    elif ext == ".docx":
        try:
            from docx import Document
            doc = Document(str(filepath))
            meta["paragraphs"] = len(doc.paragraphs)
        except Exception:
            meta["paragraphs"] = None

    # Estimate tokens from content
    try:
        converter = CONVERTERS[ext]
        content = converter(filepath)
        meta["estimated_tokens"] = len(content) // 4
    except Exception:
        meta["estimated_tokens"] = None

    return meta


# ── Frontmatter ─────────────────────────────────────────────────────────


def make_frontmatter(filepath: Path) -> str:
    return (
        "---\n"
        f"source: {filepath.name}\n"
        f"format: {filepath.suffix.lower().lstrip('.')}\n"
        f"converted_by: shuck v{__version__}\n"
        "---\n\n"
    )


# ── Converter dispatch ──────────────────────────────────────────────────


CONVERTERS = {
    ".docx": convert_docx,
    ".pdf":  convert_pdf,
    ".xlsx": convert_xlsx,
    ".pptx": convert_pptx,
    ".csv":  convert_csv,
}


# ── CLI ─────────────────────────────────────────────────────────────────


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        prog="shuck",
        description="Convert documents to clean Markdown for AI agents and LLMs.",
    )
    parser.add_argument("file", nargs="?", help="Path to the document to convert")
    parser.add_argument("-o", "--output", metavar="FILE", help="Write output to a specific file")
    parser.add_argument("-d", "--dir", metavar="DIR", help="Write output to a directory (auto-named)")
    parser.add_argument("--meta", action="store_true", help="Output JSON metadata (for agent probing)")
    parser.add_argument("--no-frontmatter", action="store_true", help="Omit YAML frontmatter from output")
    parser.add_argument("--version", action="version", version=f"shuck {__version__}")
    parser.add_argument("--formats", action="store_true", help="List supported formats and exit")
    return parser


def main():
    parser = build_parser()
    args = parser.parse_args()

    # --formats: list supported formats
    if args.formats:
        print("Supported formats:")
        for ext, desc in SUPPORTED_FORMATS.items():
            print(f"  {ext:8s} {desc}")
        sys.exit(0)

    # Require file argument for all other modes
    if not args.file:
        parser.print_help()
        sys.exit(1)

    filepath = Path(args.file).resolve()
    if not filepath.exists():
        print(f"Error: file not found: {filepath}", file=sys.stderr)
        sys.exit(1)

    ext = filepath.suffix.lower()
    if ext not in CONVERTERS:
        supported = ", ".join(SUPPORTED_FORMATS.keys())
        print(f"Error: unsupported format '{ext}'. Supported: {supported}", file=sys.stderr)
        sys.exit(1)

    # --meta mode
    if args.meta:
        meta = get_meta(filepath)
        print(json.dumps(meta, indent=2, ensure_ascii=False))
        sys.exit(0)

    # Convert
    converter = CONVERTERS[ext]
    md_content = converter(filepath)

    # Add frontmatter unless disabled
    if not args.no_frontmatter:
        md_content = make_frontmatter(filepath) + md_content

    # Output
    if args.output:
        out_path = Path(args.output)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        out_path.write_text(md_content, encoding="utf-8")
        print(f"Written to: {out_path}", file=sys.stderr)
    elif args.dir:
        out_dir = Path(args.dir)
        out_dir.mkdir(parents=True, exist_ok=True)
        out_name = filepath.stem + ".md"
        out_path = out_dir / out_name
        counter = 1
        while out_path.exists():
            out_path = out_dir / f"{filepath.stem} ({counter}).md"
            counter += 1
        out_path.write_text(md_content, encoding="utf-8")
        print(f"Written to: {out_path}", file=sys.stderr)
    else:
        # stdout (default, agent-friendly)
        sys.stdout.write(md_content)


if __name__ == "__main__":
    main()
