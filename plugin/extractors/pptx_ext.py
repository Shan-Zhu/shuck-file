"""PowerPoint (.pptx) extractor for shuck-file."""

import sys
from pathlib import Path
from .base import BaseExtractor


class PptxExtractor(BaseExtractor):

    def extract(self, filepath: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            print("Missing dependency: python-pptx\n  pip install python-pptx", file=sys.stderr)
            sys.exit(1)

        prs = Presentation(str(filepath))
        sections: list[str] = []

        for idx, slide in enumerate(prs.slides, 1):
            slide_lines = [f"## Slide {idx}"]

            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_lines.append(f"\n### {slide.shapes.title.text.strip()}\n")

            for shape in slide.shapes:
                if shape.shape_id == getattr(slide.shapes.title, "shape_id", None):
                    continue

                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_lines.append(text)

                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.replace("|", "\\|").replace("\n", " ") for cell in row.cells]
                        rows.append(cells)
                    if rows:
                        max_cols = max(len(r) for r in rows)
                        for r in rows:
                            while len(r) < max_cols:
                                r.append("")
                        slide_lines.append("")
                        slide_lines.append("| " + " | ".join(rows[0]) + " |")
                        slide_lines.append("| " + " | ".join(["---"] * max_cols) + " |")
                        for r in rows[1:]:
                            slide_lines.append("| " + " | ".join(r) + " |")
                        slide_lines.append("")

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_lines.append(f"\n> **Notes:** {notes}\n")

            sections.append("\n".join(slide_lines))

        return "\n\n".join(sections).strip() + "\n"

    def estimate_tokens(self, filepath: Path) -> int:
        return filepath.stat().st_size // 4

    def extract_tables(self, filepath: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            return ""

        prs = Presentation(str(filepath))
        tables_md = []

        for idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.replace("|", "\\|").replace("\n", " ") for cell in row.cells]
                        rows.append(cells)
                    if rows:
                        max_cols = max(len(r) for r in rows)
                        for r in rows:
                            while len(r) < max_cols:
                                r.append("")
                        md = [f"\n**Slide {idx} Table**\n"]
                        md.append("| " + " | ".join(rows[0]) + " |")
                        md.append("| " + " | ".join(["---"] * max_cols) + " |")
                        for r in rows[1:]:
                            md.append("| " + " | ".join(r) + " |")
                        tables_md.append("\n".join(md))

        return "\n\n".join(tables_md).strip() + "\n" if tables_md else ""
